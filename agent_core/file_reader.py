"""
file_reader.py
One shared "read anything" module used by both guideline ingestion and draft
ingestion, so a person can hand this agent a PDF, a Word doc, a PowerPoint,
an Excel/CSV sheet, a plain text/markdown/code file, an image, or a whole
.zip / folder mixing all of the above — and it always comes back as usable
text (or, for images, is flagged for vision-based reading by the caller).

Design goals:
  - Never require a specific extension list from the user — try structured
    readers first, then fall back to a raw-bytes-as-text attempt, so nothing
    is silently dropped just because its extension is unfamiliar.
  - Zip / folder trees are walked recursively (including nested zips),
    skipping OS junk (.DS_Store, __MACOSX, Thumbs.db, resource forks).
  - Every function degrades gracefully: a single unreadable file inside a
    batch is skipped with a reason, it never kills the whole batch.
"""

import os
import zipfile
import tempfile

TEXT_LIKE_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".yaml", ".yml",
    ".html", ".htm", ".xml", ".rtf", ".log", ".ini", ".cfg", ".conf",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".h",
    ".css", ".sql", ".sh", ".rst",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".heic"}
ARCHIVE_EXTENSIONS = {".zip"}
UNSUPPORTED_BINARY_EXTENSIONS = {
    ".exe", ".dmg", ".mp4", ".mov", ".mp3", ".wav", ".avi", ".mkv", ".dll", ".bin",
}

_APPLEDOUBLE_MAGIC = b"\x00\x05\x16\x07"
_JUNK_NAMES = {".ds_store", "thumbs.db", "desktop.ini"}


def is_junk_name(filename: str) -> bool:
    """Cheap, name-only junk check — safe to call before a file even hits disk."""
    name = os.path.basename(filename)
    lower = name.lower()
    if lower in _JUNK_NAMES or name.startswith("._") or "__MACOSX" in filename:
        return True
    return False


def is_junk_file(path: str) -> bool:
    """Fuller check once the file is on disk — also sniffs AppleDouble magic bytes,
    since some upload paths mangle the leading dot of hidden files."""
    if is_junk_name(path):
        return True
    try:
        with open(path, "rb") as f:
            if f.read(4) == _APPLEDOUBLE_MAGIC:
                return True
    except OSError:
        pass
    return False


def classify(path: str) -> str:
    """Returns one of: 'image', 'archive', 'unsupported', 'document'."""
    ext = os.path.splitext(path)[1].lower()
    if ext in IMAGE_EXTENSIONS:
        return "image"
    if ext in ARCHIVE_EXTENSIONS:
        return "archive"
    if ext in UNSUPPORTED_BINARY_EXTENSIONS:
        return "unsupported"
    return "document"


def extract_text(path: str) -> str:
    """
    Best-effort text extraction from any non-image, non-archive file.
    Tries a structured reader by extension first (pdf/docx/pptx/xlsx/json),
    then falls back to a tolerant plain-text read for everything else
    (covers .txt/.md/.csv/code files and any extension we don't special-case).
    Raises ValueError with a clear reason if nothing readable comes out.
    """
    ext = os.path.splitext(path)[1].lower()
    name = os.path.basename(path)

    if ext == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(path)
        text = "\n".join(page.extract_text() or "" for page in reader.pages).strip()
        if not text:
            raise ValueError(
                f"'{name}' looks like a scanned/image-only PDF — no extractable text. "
                "Try re-exporting it as text, or upload page screenshots as images instead."
            )
        return text

    if ext == ".docx":
        import docx
        doc = docx.Document(path)
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        text = "\n".join(parts).strip()
        if not text:
            raise ValueError(f"'{name}' has no readable text content.")
        return text

    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            lines.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        lines.append(" | ".join(cell.text for cell in row.cells))
            if lines:
                slides_text.append(f"Slide {i}: " + " / ".join(lines))
        text = "\n".join(slides_text).strip()
        if not text:
            raise ValueError(f"'{name}' has no readable text content on any slide.")
        return text

    if ext in (".xlsx", ".xlsm"):
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        parts = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True, max_row=200):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(", ".join(cells))
            if rows:
                parts.append(f"Sheet '{sheet.title}':\n" + "\n".join(rows))
        text = "\n\n".join(parts).strip()
        if not text:
            raise ValueError(f"'{name}' appears to be an empty spreadsheet.")
        return text

    if ext == ".json":
        with open(path, "r", encoding="utf-8") as f:
            return f.read()

    # Everything else (.txt/.md/.csv/code files/unrecognized extensions):
    # a tolerant plain-text read. This is deliberately permissive rather than
    # keeping an exhaustive extension allowlist, so nothing gets rejected
    # just for having an unfamiliar suffix.
    try:
        with open(path, "rb") as f:
            raw = f.read()
    except OSError as e:
        raise ValueError(f"Couldn't open '{name}': {e}")

    try:
        content = raw.decode("utf-8")
    except UnicodeDecodeError:
        try:
            content = raw.decode("latin-1")
        except Exception:
            raise ValueError(
                f"'{name}' doesn't look like a text-readable file. Supported: PDF, DOCX, "
                "PPTX, XLSX, images, plain text/code files, or a .zip/folder of these."
            )

    if not content.strip():
        raise ValueError(f"'{name}' has no readable text content.")
    return content


def read_image_bytes(path: str):
    """Returns (bytes, media_type) for an image file, normalizing exotic
    extensions down to a type Claude's vision API accepts."""
    import mimetypes
    ext = os.path.splitext(path)[1].lower()
    mt, _ = mimetypes.guess_type(path)
    if not mt or not mt.startswith("image/") or ext in (".bmp", ".tiff", ".heic"):
        mt = "image/png" if ext == ".png" else "image/jpeg"
    with open(path, "rb") as f:
        return f.read(), mt


def walk_archive(path: str):
    """
    Yields (file_path, is_archive) for every real (non-junk) file inside a
    .zip, recursing into nested zips, using a temp dir that's cleaned up by
    the caller's `with` — so this is a generator meant to be consumed inside
    a `with tempfile.TemporaryDirectory()` block. For simplicity this
    extracts fully up front and returns a flat list instead of streaming.
    """
    results = []
    with zipfile.ZipFile(path) as zf:
        tmp_dir = tempfile.mkdtemp(prefix="agent11_zip_")
        zf.extractall(tmp_dir)
    for root, _dirs, files in os.walk(tmp_dir):
        if "__MACOSX" in root:
            continue
        for name in sorted(files):
            if is_junk_name(name):
                continue
            file_path = os.path.join(root, name)
            if is_junk_file(file_path):
                continue
            results.append(file_path)
    return results


def collect_files(paths: list) -> list:
    """
    Expands a flat list of uploaded file paths into a flat list of *real*
    content files: any .zip in the input is unpacked (recursively) and its
    contents included; junk files are dropped. Order is preserved as much as
    possible. This is the single place both guideline and draft ingestion
    call to turn "whatever the user attached" into a clean file list.
    """
    expanded = []
    for path in paths:
        if is_junk_file(path):
            continue
        ext = os.path.splitext(path)[1].lower()
        if ext in ARCHIVE_EXTENSIONS:
            try:
                expanded.extend(collect_files(walk_archive(path)))
            except zipfile.BadZipFile:
                raise ValueError(f"'{os.path.basename(path)}' isn't a valid .zip archive.")
        else:
            expanded.append(path)
    return expanded
